"""
Gera slides.pptx para o Projeto Final SD (XRSC09).
Paleta: Ocean Gradient (azul profundo / teal / midnight) + branco.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

# Paleta
NAVY     = RGBColor(0x06, 0x5A, 0x82)
TEAL     = RGBColor(0x1C, 0x72, 0x93)
MIDNIGHT = RGBColor(0x21, 0x29, 0x5C)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xF4, 0xF6, 0xF8)
GRAY     = RGBColor(0x55, 0x55, 0x55)
ACCENT   = RGBColor(0xFF, 0xC1, 0x07)
DARK     = RGBColor(0x1A, 0x1A, 0x2A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

DIAG_DIR = os.path.join(os.path.dirname(__file__), "..", "diagrams")

def add_bg(slide, color=LIGHT):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def add_band(slide, top, height, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SW, height)
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tx

def add_bullets(slide, x, y, w, h, items, size=16, color=DARK, bold_first_word=False):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tx

def add_footer(slide, n):
    # Barra de rodapé
    add_band(slide, SH - Inches(0.35), Inches(0.35), NAVY)
    add_text(slide, Inches(0.3), SH - Inches(0.32), Inches(12), Inches(0.3),
             "XRSC09 — Sistemas Distribuídos  •  Grupo XX  •  Projeto Final", 10, color=WHITE)
    add_text(slide, SW - Inches(1.0), SH - Inches(0.32), Inches(0.7), Inches(0.3),
             f"{n}", 10, color=WHITE, align=PP_ALIGN.RIGHT, bold=True)

def add_header(slide, title, subtitle=None):
    # Faixa superior
    add_band(slide, 0, Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6), title,
             size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
                 subtitle, size=13, color=TEAL, bold=True)

def slide_blank(n_footer=None, bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, bg)
    if n_footer is not None:
        add_footer(s, n_footer)
    return s

# ===================================================================
# SLIDE 1 — Capa
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, MIDNIGHT)
# Faixa de cor de destaque
add_band(s, Inches(2.4), Inches(0.06), ACCENT)
add_text(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.0),
         "Plataforma Distribuída", size=42, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.4),
         "para Acompanhamento de Manutenção Veicular", size=32, bold=True, color=ACCENT)
add_text(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.6),
         "Arquitetura orientada a eventos com RabbitMQ • Edge + Cloud", size=20, color=WHITE)
add_text(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.4),
         "XRSC09 — Sistemas Distribuídos  •  Projeto Final", size=16, color=TEAL, bold=True)
add_text(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4),
         "Grupo XX — Junho de 2026", size=14, color=WHITE)
add_text(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.4),
         "Membro 1  •  Membro 2  •  Membro 3  •  Membro 4  •  Membro 5",
         size=13, color=WHITE)

# ===================================================================
# SLIDE 2 — Problema
# ===================================================================
s = slide_blank(2)
add_header(s, "O problema", "Transparência na manutenção veicular")

add_text(s, Inches(0.6), Inches(1.5), Inches(7.5), Inches(0.5),
         "Cenário sorteado:", 18, bold=True, color=NAVY)
add_text(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(2.4),
         "Empresa de manutenção de veículos quer permitir que o cliente acompanhe "
         "o serviço em tempo real, com etapas como relato, diagnóstico, identificação "
         "da causa, execução do reparo e substituição rastreável de peças.",
         14, color=DARK)

# Dor
add_band(s, Inches(4.2), Inches(0.05), TEAL)
add_text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.5),
         "Dor central", 18, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(4.9), Inches(12), Inches(2.2), [
    "Cliente não vê o que foi feito → desconfiança e reclamações",
    "Sem evidência verificável de cada etapa do serviço",
    "Peças substituídas sem rastreio claro",
    "Comunicação informal e fragmentada entre cliente e oficina"
], size=14)

# ===================================================================
# SLIDE 3 — Motivação para abordagem distribuída
# ===================================================================
s = slide_blank(3)
add_header(s, "Por que uma solução distribuída?", "Limites do monolito")

cards = [
    ("Escala desigual", "Processar vídeo escala diferente de servir REST. Tudo no mesmo processo amarra as cargas."),
    ("Vários destinatários", "Mudar etapa gera notificação, auditoria, dashboard e e-mail — fan-out natural."),
    ("Resiliência geográfica", "A oficina precisa registrar evidência mesmo offline e sincronizar depois."),
    ("Extensibilidade", "Novos consumers (billing, relatório, seguradora) devem entrar sem mexer no núcleo."),
]

for i, (title, body) in enumerate(cards):
    col, row = i % 2, i // 2
    x = Inches(0.6 + col * 6.3)
    y = Inches(1.6 + row * 2.5)
    w, h = Inches(6.0), Inches(2.3)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = TEAL; card.line.width = Pt(1.5)
    add_text(s, x + Inches(0.3), y + Inches(0.2), w - Inches(0.6), Inches(0.5),
             title, 18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y + Inches(0.8), w - Inches(0.6), h - Inches(1.0),
             body, 13, color=DARK)

# ===================================================================
# SLIDE 4 — Planta da oficina
# ===================================================================
s = slide_blank(4)
add_header(s, "Planta da oficina", "Onde ficam câmeras, terminais e o gateway local")
s.shapes.add_picture(os.path.join(DIAG_DIR, "planta_oficina.png"),
                     Inches(1.0), Inches(1.2), height=Inches(5.7))
add_text(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
         "4 galpões com câmera IP + tablet • almoxarifado • recepção • sala técnica com gateway, NAS e mini-broker",
         11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 5 — C4 Nível 1 Contexto
# ===================================================================
s = slide_blank(5)
add_header(s, "Modelo C4 — Nível 1: Contexto", "Atores e sistemas externos")
s.shapes.add_picture(os.path.join(DIAG_DIR, "c4_nivel1_contexto.png"),
                     Inches(1.0), Inches(1.2), height=Inches(5.5))
add_text(s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
         "5 atores • integração com câmeras, storage, gateways, ERP, pagamento, NF-e",
         11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 6 — C4 Nível 2 Containers
# ===================================================================
s = slide_blank(6)
add_header(s, "Modelo C4 — Nível 2: Containers", "Zona de borda (oficina) + zona de nuvem")
s.shapes.add_picture(os.path.join(DIAG_DIR, "c4_nivel2_containers.png"),
                     Inches(1.0), Inches(1.2), height=Inches(5.6))
add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
         "Edge: gateway, NAS, mini-broker federado    •    Cloud: API, Postgres, RabbitMQ cluster, workers, storage",
         11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 7 — C4 Nível 3 Backend
# ===================================================================
s = slide_blank(7)
add_header(s, "C4 Nível 3 — Componentes do Backend", "Controllers • Services • Repositories • EventPublisher")
s.shapes.add_picture(os.path.join(DIAG_DIR, "c4_nivel3_componentes.png"),
                     Inches(1.0), Inches(1.2), height=Inches(5.6))
add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
         "EventPublisher é o ponto único de saída para o broker — substituir middleware afeta só essa classe",
         11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 8 — RabbitMQ topology
# ===================================================================
s = slide_blank(8)
add_header(s, "Topologia RabbitMQ", "5 exchanges • ~12 filas • ~25 routing keys planejadas")
s.shapes.add_picture(os.path.join(DIAG_DIR, "rabbitmq_topology.png"),
                     Inches(1.0), Inches(1.2), height=Inches(5.6))
add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
         "Topic + Fanout + Direct + Headers + DLX  • patrões Pub/Sub, Work Queue, Topic Routing, Headers, DLX",
         11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 9 — Exchanges detalhadas
# ===================================================================
s = slide_blank(9)
add_header(s, "Cinco exchanges, cinco padrões", "Cada tipo de AMQP cobre um caso de uso")

exchanges = [
    ("oficina.events",       "Topic",   "Hub principal — routing key dominio.entidade.acao",      "#F57F17"),
    ("oficina.broadcast",    "Fanout",  "Avisos globais (manutenção da plataforma, alarmes)",     "#AD1457"),
    ("oficina.commands",     "Direct",  "Comandos diretos para um worker específico",             "#283593"),
    ("oficina.notifications","Headers", "Roteia notificação por canal (email/sms/whatsapp)",      "#4527A0"),
    ("oficina.dlx",          "DLX",     "Dead Letter — mensagens com falha para inspeção/retry", "#B71C1C"),
]
y0 = Inches(1.6)
for i, (name, typ, desc, color) in enumerate(exchanges):
    y = y0 + Inches(i * 1.0)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.6), y, Inches(2.6), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
    box.line.fill.background()
    add_text(s, Inches(0.7), y + Inches(0.10), Inches(2.4), Inches(0.4),
             name, 13, bold=True, color=WHITE)
    add_text(s, Inches(0.7), y + Inches(0.45), Inches(2.4), Inches(0.35),
             typ, 11, color=WHITE)
    add_text(s, Inches(3.4), y + Inches(0.20), Inches(9.5), Inches(0.6),
             desc, 14, color=DARK)

# ===================================================================
# SLIDE 10 — Máquina de estados
# ===================================================================
s = slide_blank(10)
add_header(s, "Máquina de estados das etapas", "12 estados, transições validadas, evento por transição")
s.shapes.add_picture(os.path.join(DIAG_DIR, "state_machine.png"),
                     Inches(2.8), Inches(1.1), height=Inches(5.9))
add_text(s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
         "Cada transição publica maintenance.step.updated no Topic Exchange oficina.events",
         11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 11 — Sequence: mídia
# ===================================================================
s = slide_blank(11)
add_header(s, "Fluxo 1 — Upload de mídia", "Resposta síncrona rápida • processamento assíncrono")
s.shapes.add_picture(os.path.join(DIAG_DIR, "sequence_midia.png"),
                     Inches(0.8), Inches(1.2), width=Inches(11.7))
add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
         "Backend responde em <100ms • media-worker processa em ~2s • cliente vê quando media.processed chega",
         11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 12 — Sequence: cascata
# ===================================================================
s = slide_blank(12)
add_header(s, "Fluxo 2 — Mudança de etapa", "Uma chamada → cascata de workers em paralelo")
s.shapes.add_picture(os.path.join(DIAG_DIR, "sequence_etapa.png"),
                     Inches(0.5), Inches(1.2), width=Inches(12.3))
add_text(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
         "notification-worker e audit-worker reagem ao mesmo evento • inventory ignora • cliente vê notificação",
         11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 13 — Deployment híbrido
# ===================================================================
s = slide_blank(13)
add_header(s, "Implantação híbrida edge + cloud", "Resiliência operacional como decisão arquitetural")
s.shapes.add_picture(os.path.join(DIAG_DIR, "deployment.png"),
                     Inches(0.8), Inches(1.2), width=Inches(11.7))
add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
         "K8s na nuvem (HPA por worker) • gateway na oficina (NAS, mini-broker federado) • funciona offline",
         11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 14 — Stack do protótipo
# ===================================================================
s = slide_blank(14)
add_header(s, "Stack do protótipo", "Tudo conteinerizado, sobe com um único comando")

stacks = [
    ("Frontend",  "Next.js 14 + React 18",       NAVY),
    ("Backend",   "Node.js 20 + Express 4",      TEAL),
    ("Broker",    "RabbitMQ 3 (management UI)",  MIDNIGHT),
    ("Banco",     "PostgreSQL 16",               NAVY),
    ("Workers",   "4× Node + amqplib",           TEAL),
    ("Orquestração", "Docker Compose",            MIDNIGHT),
]
cols = 3
for i, (k, v, color) in enumerate(stacks):
    col, row = i % cols, i // cols
    x = Inches(0.6 + col * 4.25)
    y = Inches(1.6 + row * 2.4)
    w, h = Inches(4.0), Inches(2.1)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.15), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.15), y, w - Inches(0.15), h)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color; card.line.width = Pt(0.5)
    add_text(s, x + Inches(0.4), y + Inches(0.3), w - Inches(0.5), Inches(0.5),
             k, 18, bold=True, color=color)
    add_text(s, x + Inches(0.4), y + Inches(0.95), w - Inches(0.5), Inches(0.8),
             v, 14, color=DARK)

add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.4),
         "$ docker compose up --build", 14, color=NAVY, bold=True, font="Consolas")

# ===================================================================
# SLIDE 15 — O que o protótipo demonstra
# ===================================================================
s = slide_blank(15)
add_header(s, "O que o protótipo demonstra", "Fatia mínima viável da arquitetura completa")

add_text(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.5),
         "Implementado no código:", 18, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5), [
    "5 exchanges (Topic, Fanout, Direct, Headers, DLX)",
    "Filas com bindings por routing key, # e *",
    "Máquina de estados de 12 etapas, transições validadas",
    "Upload de mídia + processamento assíncrono",
    "Rastreio de peças (request → reserve → install)",
    "Notificações personalizadas para stakeholders",
    "Trilha de auditoria com binding universal #",
    "Frontend autenticado por papel (4 perfis)"
], size=14)

add_text(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(0.5),
         "Defendido na apresentação:", 18, bold=True, color=NAVY)
add_bullets(s, Inches(7.0), Inches(2.0), Inches(6.0), Inches(4.5), [
    "Arquitetura edge-cloud com mini-broker federado",
    "Gateway local com NAS para vídeos brutos",
    "Cluster RabbitMQ com quorum queues em K8s",
    "HPA por worker, observabilidade Prometheus/Grafana",
    "Integração real com email/SMS/WhatsApp (Headers)",
    "Storage S3/R2 para vídeos definitivos",
    "Report worker, billing worker, NF-e",
    "Push em tempo real via WebSocket"
], size=14)

# ===================================================================
# SLIDE 16 — Resultados dos cenários
# ===================================================================
s = slide_blank(16)
add_header(s, "Cenários executados", "Quatro fluxos end-to-end com a stack levantada")

scenarios = [
    ("C1 — Mudança de etapa",      "Backend → maintenance.step.updated → notification + audit + portal",
     "<500 ms"),
    ("C2 — Upload de mídia",       "3 arquivos → media.uploaded → worker → media.processed → notif",
     "~7 s p/ 3 arquivos"),
    ("C3 — Solicitação de peça",   "parts.requested → inventory-worker → reserve → parts.reserved",
     "~1 s"),
    ("C4 — Auditoria universal",   "audit-worker assina '#' → 11 eventos capturados nos cenários 1+2+3",
     "100% cobertura"),
]
y0 = Inches(1.6)
for i, (title, body, metric) in enumerate(scenarios):
    y = y0 + Inches(i * 1.25)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.15), Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), y, Inches(12.0), Inches(1.05))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = TEAL; card.line.width = Pt(0.3)
    add_text(s, Inches(0.85), y + Inches(0.10), Inches(7.0), Inches(0.4),
             title, 15, bold=True, color=NAVY)
    add_text(s, Inches(0.85), y + Inches(0.5), Inches(8.5), Inches(0.5),
             body, 12, color=DARK)
    add_text(s, Inches(9.5), y + Inches(0.3), Inches(3.0), Inches(0.5),
             metric, 18, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)

# ===================================================================
# SLIDE 17 — Limitações e trabalhos futuros
# ===================================================================
s = slide_blank(17)
add_header(s, "Limitações e trabalhos futuros", "O que o protótipo não cobre — e como evoluir")

add_text(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.5),
         "Limitações conhecidas:", 18, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5), [
    "Sessão em memória (impede escala do backend)",
    "Processamento de mídia é simulado (sem FFmpeg)",
    "Edge gateway só na apresentação",
    "Sem WebSocket — cliente faz polling manual",
    "Sem integração real com email/SMS/WhatsApp",
    "Sem upload para S3/R2"
], size=14)

add_text(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(0.5),
         "Trabalhos futuros:", 18, bold=True, color=NAVY)
add_bullets(s, Inches(7.0), Inches(2.0), Inches(6.0), Inches(4.5), [
    "Implementar edge gateway com federation",
    "Quorum queues no broker",
    "Push em tempo real via WebSocket / SSE",
    "Integração real com gateways via Headers Exchange",
    "FFmpeg para thumbnail + transcodificação",
    "Report worker (PDF da OS concluída)",
    "Prometheus + Grafana para observabilidade",
    "Integração com NF-e e billing real"
], size=13)

# ===================================================================
# SLIDE 18 — Demo ao vivo (slide marcador)
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, MIDNIGHT)
add_text(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.5),
         "Demo ao vivo", size=60, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.6),
         "docker compose up", size=24, color=ACCENT, font="Consolas")
add_bullets(s, Inches(2.0), Inches(4.8), Inches(9.5), Inches(2.0), [
    "frontend  http://localhost:3000",
    "backend   http://localhost:4000/health",
    "RabbitMQ Management  http://localhost:15672  (guest/guest)"
], size=16, color=WHITE)
add_footer(s, 18)

# ===================================================================
# SLIDE 19 — Conclusão
# ===================================================================
s = slide_blank(19)
add_header(s, "Conclusão", "Síntese e contribuições")

add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
         "O que entregamos:", 18, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(2.0), [
    "Projeto completo de sistema distribuído orientado a eventos para um problema real e atual",
    "Modelo C4 nos 4 níveis + planta + máquina de estados + 2 sequence diagrams + deployment",
    "Protótipo funcional com 5 exchanges, 4 workers, 9 tabelas, frontend autenticado, Docker Compose"
], size=14)

add_text(s, Inches(0.6), Inches(4.1), Inches(12), Inches(0.5),
         "Lições aprendidas:", 18, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(4.6), Inches(12.0), Inches(2.5), [
    "Desacoplar via broker é o que viabiliza escalar workers de forma independente",
    "A trilha de auditoria 'de graça' pelo padrão # é uma vantagem técnica concreta do Topic Exchange",
    "A arquitetura edge-cloud não é luxo: é o que mantém a oficina operacional durante quedas",
    "Centralizar a publicação em EventPublisher isola a decisão de middleware"
], size=14)

# ===================================================================
# SLIDE 20 — Perguntas
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, MIDNIGHT)
add_band(s, Inches(3.3), Inches(0.06), ACCENT)
add_text(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.5),
         "Obrigado.", size=72, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.5),
         "Perguntas?", size=28, color=ACCENT)
add_text(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.4),
         "github.com/<usuario>/projeto-final-sd", size=14, color=TEAL, font="Consolas")
add_footer(s, 20)

# Save
out = os.path.join(os.path.dirname(__file__), "slides.pptx")
prs.save(out)
print(f"OK: {out} ({len(prs.slides)} slides)")
